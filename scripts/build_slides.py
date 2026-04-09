#!/usr/bin/env python3
"""Build a pptx presentation from markdown + template.

Usage:
    python3 scripts/build_slides.py input.md [-t template.pptx] [-o output.pptx] [--img-dir slides-img]

The markdown format uses YAML front matter for metadata, ## headings for slide
titles, and --- separators between slides. See templates/template.md for format
documentation.
"""

from __future__ import annotations

import argparse
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu

from slide_utils import (
    PresentationData, SlideData, TextElement, BlockquoteElement,
    TableElement, ImageElement, MermaidElement, ChartElement,
    Paragraph, Run,
    parse_markdown, find_layout, find_layout_flexible, select_layout,
    set_text_frame, collect_text_paragraphs,
    add_table_to_slide, add_image_to_slide,
    render_all_mermaid, wps_fixup,
)
from chart_utils import add_chart_to_slide, parse_width_pct

# Default content-area geometry (EMU) — matches the slide master content region
_CL  = Emu(838200)    # content left   (~0.69")
_CT  = Emu(1825625)   # content top    (~1.50")
_CW  = Emu(10515600)  # content width  (~8.67")
_CH  = Emu(4500000)   # content height (~3.71")
_GAP = Emu(200000)    # column gap     (~0.165")


# ---------------------------------------------------------------------------
# Slide population
# ---------------------------------------------------------------------------

def _set_title(slide, title: str):
    """Set the title placeholder text. Uses slide.shapes.title when available."""
    shape = slide.shapes.title
    if shape and shape.has_text_frame:
        shape.text = title
        return
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0 and ph.has_text_frame:
            ph.text = title
            return


def _stack_images(slide, image_elements, left, top, max_w, max_h):
    """Stack images vertically within the given bounding box."""
    n = len(image_elements)
    per_h = max_h // n
    for ii, img_el in enumerate(image_elements):
        add_image_to_slide(slide, img_el.path, left, top + per_h * ii, max_w, per_h)


def populate_title_slide(slide, slide_data: SlideData, pdata: PresentationData):
    """Populate a Title Slide layout (idx 0=title, 1=subtitle)."""
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            ph.text = pdata.title
        elif idx == 1:
            lines = pdata.subtitle.split('\n')
            ph.text_frame.clear()
            for i, line in enumerate(lines):
                if i == 0:
                    ph.text_frame.paragraphs[0].text = line.strip()
                else:
                    ph.text_frame.add_paragraph().text = line.strip()


def populate_standard_layout(slide, slide_data: SlideData):
    """Populate a Title and Content layout (idx 0=title, 1=body)."""
    _set_title(slide, slide_data.title)

    text_elements = [e for e in slide_data.body_elements
                     if isinstance(e, (TextElement, BlockquoteElement))]
    table_elements = [e for e in slide_data.body_elements if isinstance(e, TableElement)]
    image_elements = [e for e in slide_data.body_elements if isinstance(e, ImageElement)]
    chart_elements = [e for e in slide_data.body_elements if isinstance(e, ChartElement)]

    if chart_elements:
        print(f"  Warning: slide \"{slide_data.title}\" 使用 standard 布局但包含图表，"
              "图表将被忽略。请改用 <!-- layout: chart --> 以正确渲染图表。")

    body_ph = next((ph for ph in slide.placeholders if ph.placeholder_format.idx == 1), None)

    if body_ph and text_elements:
        all_paras = collect_text_paragraphs(text_elements)
        if all_paras:
            set_text_frame(body_ph.text_frame, all_paras)

    if table_elements:
        if body_ph:
            tbl_left = body_ph.left
            tbl_top = body_ph.top + body_ph.height // 2 if text_elements else body_ph.top
            tbl_width = body_ph.width
            tbl_height = body_ph.height // 2 if text_elements else body_ph.height
        else:
            tbl_left, tbl_top = _CL, _CT
            tbl_width, tbl_height = _CW, Emu(2500000)
        for tbl_el in table_elements:
            add_table_to_slide(slide, tbl_el, tbl_left, tbl_top, tbl_width, tbl_height)

    if image_elements:
        if body_ph:
            img_left, img_top = body_ph.left, body_ph.top
            img_max_w, img_max_h = body_ph.width, body_ph.height
        else:
            img_left, img_top = _CL, _CT
            img_max_w, img_max_h = _CW, _CH
        _stack_images(slide, image_elements, img_left, img_top, img_max_w, img_max_h)


def populate_title_only(slide, slide_data: SlideData):
    """Title Only layout — title + free-form images."""
    _set_title(slide, slide_data.title)
    image_elements = [e for e in slide_data.body_elements if isinstance(e, ImageElement)]
    if image_elements:
        _stack_images(slide, image_elements,
                      _CL, _CT, _CW, _CH)


def populate_chart_layout(slide, slide_data: SlideData):
    """Title Only layout with free-form charts. Pure chart rendering — no text."""
    _set_title(slide, slide_data.title)

    chart_elements = [e for e in slide_data.body_elements if isinstance(e, ChartElement)]
    if not chart_elements:
        return

    if len(chart_elements) == 1:
        chart_el = chart_elements[0]
        chart_w = int(_CW * parse_width_pct(chart_el.width))
        if chart_el.position == "left":
            chart_left = _CL
        elif chart_el.position == "right":
            chart_left = _CL + _CW - chart_w
        else:
            chart_left = _CL + (_CW - chart_w) // 2
        add_chart_to_slide(slide, chart_el, chart_left, _CT, chart_w, _CH)
    else:
        chart_w = (_CW - _GAP) // 2
        rows = (len(chart_elements) + 1) // 2
        chart_h = (_CH - _GAP * (rows - 1)) // rows
        for ci, chart_el in enumerate(chart_elements):
            row, col = divmod(ci, 2)
            add_chart_to_slide(slide, chart_el,
                               _CL + col * (chart_w + _GAP),
                               _CT + row * (chart_h + _GAP),
                               chart_w, chart_h)


def _two_col_region(slide):
    """Return (left_ph, right_ph, left, top, width, height) for two-column layouts."""
    left_ph = right_ph = None
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 1:
            left_ph = ph
        elif idx == 2:
            right_ph = ph

    if left_ph:
        left   = left_ph.left
        top    = left_ph.top
        height = left_ph.height
        width  = (right_ph.left + right_ph.width if right_ph else left_ph.left + left_ph.width) - left
    else:
        left, top = _CL, _CT
        width, height = _CW, _CH

    return left_ph, right_ph, left, top, width, height


def populate_two_content_layout(slide, slide_data: SlideData):
    """Two Content layout — blockquote → left placeholder, text → right placeholder."""
    _set_title(slide, slide_data.title)

    bq_elements    = [e for e in slide_data.body_elements if isinstance(e, BlockquoteElement)]
    text_elements  = [e for e in slide_data.body_elements if isinstance(e, TextElement)]
    table_elements = [e for e in slide_data.body_elements if isinstance(e, TableElement)]
    image_elements = [e for e in slide_data.body_elements if isinstance(e, ImageElement)]

    left_ph, right_ph, rl, rt, rw, rh = _two_col_region(slide)
    half_w = (rw - _GAP) // 2
    right_l = rl + half_w + _GAP

    if bq_elements:
        left_paras  = [p for bq in bq_elements for p in bq.paragraphs]
        right_paras = collect_text_paragraphs(text_elements)
    else:
        mid = max(1, len(text_elements) // 2)
        left_paras  = collect_text_paragraphs(text_elements[:mid])
        right_paras = collect_text_paragraphs(text_elements[mid:])

    if left_paras:
        if left_ph:
            set_text_frame(left_ph.text_frame, left_paras)
        else:
            set_text_frame(slide.shapes.add_textbox(rl, rt, half_w, rh).text_frame, left_paras)

    if right_paras:
        if right_ph:
            set_text_frame(right_ph.text_frame, right_paras)
        else:
            set_text_frame(slide.shapes.add_textbox(right_l, rt, half_w, rh).text_frame, right_paras)

    for tbl_el in table_elements:
        add_table_to_slide(slide, tbl_el, right_l, rt, half_w, rh)

    if image_elements:
        _stack_images(slide, image_elements, right_l, rt, half_w, rh)


def populate_mixed_layout(slide, slide_data: SlideData):
    """Mixed layout — text/table left, chart(s) right."""
    _set_title(slide, slide_data.title)

    text_elements  = [e for e in slide_data.body_elements
                      if isinstance(e, (TextElement, BlockquoteElement))]
    table_elements = [e for e in slide_data.body_elements if isinstance(e, TableElement)]
    chart_elements = [e for e in slide_data.body_elements if isinstance(e, ChartElement)]

    left_ph, right_ph, rl, rt, rw, rh = _two_col_region(slide)
    half_w = (rw - _GAP) // 2
    right_l = rl + half_w + _GAP

    # Left: text
    all_text_paras = []
    for e in text_elements:
        if isinstance(e, BlockquoteElement):
            all_text_paras.extend(e.paragraphs)
        else:
            all_text_paras.extend(collect_text_paragraphs([e]))

    if all_text_paras:
        if left_ph:
            set_text_frame(left_ph.text_frame, all_text_paras)
        else:
            set_text_frame(slide.shapes.add_textbox(rl, rt, half_w, rh).text_frame, all_text_paras)

    # Left: table (below text if any)
    for tbl_el in table_elements:
        tbl_top = rt + (rh // 2 if all_text_paras else 0)
        tbl_h   = rh // 2 if all_text_paras else rh
        add_table_to_slide(slide, tbl_el, rl, tbl_top, half_w, tbl_h)

    # Right: charts (use right placeholder geometry if available)
    if chart_elements:
        chart_left = right_ph.left if right_ph else right_l
        chart_top  = right_ph.top  if right_ph else rt
        chart_w    = right_ph.width  if right_ph else half_w
        chart_h    = right_ph.height if right_ph else rh
        n = len(chart_elements)
        per_h = (chart_h - _GAP * (n - 1)) // n
        for ci, chart_el in enumerate(chart_elements):
            add_chart_to_slide(slide, chart_el,
                               chart_left, chart_top + ci * (per_h + _GAP),
                               chart_w, per_h)


def populate_section_layout(slide, slide_data: SlideData):
    """Section Header layout — title (idx=0) + short description (idx=1)."""
    _set_title(slide, slide_data.title)
    text_elements = [e for e in slide_data.body_elements if isinstance(e, TextElement)]
    paras = collect_text_paragraphs(text_elements)
    body_ph = next((ph for ph in slide.placeholders if ph.placeholder_format.idx == 1), None)
    if body_ph and paras:
        set_text_frame(body_ph.text_frame, paras)


def populate_toc_layout(slide, slide_data: SlideData):
    """Populate a TOC layout: idx 0 = title, idx 1..N = one chapter item each."""
    toc_items = [
        "".join(r.text for r in para.runs).strip()
        for el in slide_data.body_elements if isinstance(el, TextElement)
        for para in el.paragraphs
        if "".join(r.text for r in para.runs).strip()
    ]

    item_phs = []
    for ph in sorted(slide.placeholders, key=lambda p: p.placeholder_format.idx):
        if ph.placeholder_format.idx == 0:
            ph.text = slide_data.title
        else:
            item_phs.append(ph)

    if not toc_items or not item_phs:
        return

    if len(item_phs) == 1:
        paras = [Paragraph(runs=[Run(text=item)], level=0) for item in toc_items]
        set_text_frame(item_phs[0].text_frame, paras)
    else:
        for i, ph in enumerate(item_phs):
            ph.text = toc_items[i] if i < len(toc_items) else ""


def populate_summary_layout(slide, slide_data: SlideData):
    populate_standard_layout(slide, slide_data)


# ---------------------------------------------------------------------------
# Main build pipeline
# ---------------------------------------------------------------------------

def build_presentation(pdata: PresentationData, template_path: str, output_path: str,
                       img_dir: str = "slides-img"):
    """Build a pptx from parsed markdown data and a template."""
    prs = Presentation(template_path)

    # Phase 1: Render mermaid diagrams
    render_all_mermaid(pdata.slides, img_dir)

    # Phase 2: Build title slide from front matter
    title_layout, _ = find_layout_flexible(prs, "cover")
    title_slide = prs.slides.add_slide(title_layout)
    populate_title_slide(title_slide, SlideData(), pdata)
    print(f"  Slide 1: \"{pdata.title}\" -> {title_layout.name}")

    # Phase 3: Build content slides
    for si, slide_data in enumerate(pdata.slides):
        layout, std_name = select_layout(slide_data, prs)
        slide = prs.slides.add_slide(layout)

        # 根据标准布局类型（std_name）路由到对应的填充函数
        if std_name == 'cover':
            populate_title_slide(slide, slide_data, pdata)
        elif std_name == 'toc':
            populate_toc_layout(slide, slide_data)
        elif std_name == 'section':
            populate_section_layout(slide, slide_data)
        elif std_name == 'summary':
            populate_summary_layout(slide, slide_data)
        elif std_name == 'two-column':
            populate_two_content_layout(slide, slide_data)
        elif std_name == 'mixed':
            populate_mixed_layout(slide, slide_data)
        elif std_name == 'chart':
            populate_chart_layout(slide, slide_data)
        elif std_name in ('image', 'title-only'):
            populate_title_only(slide, slide_data)
        else:  # standard, table, mixed 及其他
            populate_standard_layout(slide, slide_data)

        print(f"  Slide {si+2}: \"{slide_data.title}\" -> {layout.name} ({std_name})")

    # Phase 4: Save
    prs.save(output_path)

    # Phase 5: WPS compatibility fixup
    wps_fixup(output_path)

    return len(pdata.slides) + 1  # +1 for title slide


def main():
    parser = argparse.ArgumentParser(description='Build pptx from markdown')
    parser.add_argument('input', help='Input markdown file')
    parser.add_argument('-t', '--template', default=None,
                        help='Template pptx (overrides front matter)')
    parser.add_argument('-o', '--output', default=None,
                        help='Output pptx path (overrides front matter)')
    parser.add_argument('--img-dir', default='slides-img',
                        help='Directory for mermaid renders (default: slides-img/)')
    args = parser.parse_args()

    if not os.path.exists(args.input):
        print(f"Error: {args.input} not found", file=sys.stderr)
        sys.exit(1)

    print(f"==> Parsing {args.input}...")
    pdata = parse_markdown(args.input)
    print(f"  Title: {pdata.title}")
    print(f"  Slides: {len(pdata.slides)}")

    template = args.template or pdata.template
    output = args.output or pdata.output

    # Resolve template path relative to the input md file's directory
    md_dir = os.path.dirname(os.path.abspath(args.input))
    if template and not os.path.isabs(template):
        template_rel = os.path.join(md_dir, template)
        if os.path.exists(template_rel):
            template = template_rel

    if not os.path.exists(template):
        # 尝试使用默认模板（相对于脚本所在目录）
        default_template = os.path.join(
            os.path.dirname(os.path.abspath(__file__)),
            '..', 'examples', 'default.pptx'
        )
        if os.path.exists(default_template):
            print(f"Warning: 模板 {template!r} 不存在，使用默认模板: {default_template}")
            template = default_template
        else:
            print(f"Error: Template {template} not found", file=sys.stderr)
            sys.exit(1)

    print(f"==> Building slides with template: {template}")
    n_slides = build_presentation(pdata, template, output, args.img_dir)

    print(f"==> Done: {output} ({n_slides} slides)")


if __name__ == '__main__':
    main()
