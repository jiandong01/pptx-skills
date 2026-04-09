"""Chart creation utilities for pptx slide generation.

Supports:
  Category charts (CategoryChartData):
    column, column-stacked, column-stacked-100
    bar, bar-stacked, bar-stacked-100
    line, line-stacked, line-stacked-100
    area, area-stacked, area-stacked-100
    pie, pie-exploded, doughnut, doughnut-exploded
    radar, radar-filled, radar-markers
  XY charts (XyChartData):
    scatter
  Bubble charts (BubbleChartData):
    bubble
  Composite charts (XML surgery):
    combo    — per-series subtype: bar | line, axis: primary | secondary
    waterfall — stacked-column illusion with transparent base
"""

from __future__ import annotations

import copy
from lxml import etree

from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.util import Pt, Emu
from pptx.oxml.ns import qn


# ---------------------------------------------------------------------------
# Chart type map (CategoryChartData types)
# ---------------------------------------------------------------------------

CATEGORY_CHART_MAP = {
    "column":             XL_CHART_TYPE.COLUMN_CLUSTERED,
    "column-stacked":     XL_CHART_TYPE.COLUMN_STACKED,
    "column-stacked-100": XL_CHART_TYPE.COLUMN_STACKED_100,
    "bar":                XL_CHART_TYPE.BAR_CLUSTERED,
    "bar-stacked":        XL_CHART_TYPE.BAR_STACKED,
    "bar-stacked-100":    XL_CHART_TYPE.BAR_STACKED_100,
    "line":               XL_CHART_TYPE.LINE_MARKERS,
    "line-stacked":       XL_CHART_TYPE.LINE_MARKERS_STACKED,
    "line-stacked-100":   XL_CHART_TYPE.LINE_MARKERS_STACKED_100,
    "area":               XL_CHART_TYPE.AREA,
    "area-stacked":       XL_CHART_TYPE.AREA_STACKED,
    "area-stacked-100":   XL_CHART_TYPE.AREA_STACKED_100,
    "pie":                XL_CHART_TYPE.PIE,
    "pie-exploded":       XL_CHART_TYPE.PIE_EXPLODED,
    "doughnut":           XL_CHART_TYPE.DOUGHNUT,
    "doughnut-exploded":  XL_CHART_TYPE.DOUGHNUT_EXPLODED,
    "radar":              XL_CHART_TYPE.RADAR_MARKERS,
    "radar-filled":       XL_CHART_TYPE.RADAR_FILLED,
}

# Label position hints per chart family
_LABEL_POS = {
    "column": XL_LABEL_POSITION.OUTSIDE_END,
    "bar":    XL_LABEL_POSITION.OUTSIDE_END,
    "line":   XL_LABEL_POSITION.ABOVE,
    "area":   XL_LABEL_POSITION.CENTER,
    "pie":    XL_LABEL_POSITION.BEST_FIT,
    "doughnut": XL_LABEL_POSITION.CENTER,
    "radar":  XL_LABEL_POSITION.ABOVE,
}

# Waterfall defaults
_WF_GAIN_COLOR  = RGBColor(0x70, 0xAD, 0x47)   # green
_WF_LOSS_COLOR  = RGBColor(0xFF, 0x00, 0x00)   # red
_WF_TOTAL_COLOR = RGBColor(0x44, 0x72, 0xC4)   # blue
_WF_BASE_COLOR  = None                           # transparent


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _parse_color(color_str: str) -> RGBColor | None:
    if not color_str:
        return None
    hex_str = color_str.lstrip('#')
    if len(hex_str) != 6:
        return None
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def parse_width_pct(width_str: str) -> float:
    s = width_str.strip().rstrip('%')
    try:
        return float(s) / 100.0
    except ValueError:
        return 0.6


def _family(chart_type: str) -> str:
    """Return base family for label-position lookup."""
    for fam in ("column", "bar", "line", "area", "pie", "doughnut", "radar"):
        if chart_type.startswith(fam):
            return fam
    return "column"


def _apply_data_labels(plot, chart_type: str, number_format: str):
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.size = Pt(10)
    if number_format:
        dl.number_format = number_format
    fam = _family(chart_type)
    pos = _LABEL_POS.get(fam)
    if pos is not None:
        try:
            dl.label_position = pos
        except Exception:
            pass


def _apply_series_colors(plot, series_list):
    for i, sd in enumerate(series_list):
        color = _parse_color(sd.color)
        if color and i < len(plot.series):
            fill = plot.series[i].format.fill
            fill.solid()
            fill.fore_color.rgb = color


def _set_chart_title(chart, title: str):
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
        for para in chart.chart_title.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(13)
                run.font.bold = True
    else:
        chart.has_title = False


def _position_chart(chart_el, content_left, content_top, content_width, content_height):
    """Compute (left, top, w, h) for a single chart given position/width hints."""
    pct = parse_width_pct(chart_el.width)
    chart_w = int(content_width * pct)
    if chart_el.position == "left":
        chart_left = content_left
    elif chart_el.position == "right":
        chart_left = content_left + content_width - chart_w
    else:
        chart_left = content_left + (content_width - chart_w) // 2
    return chart_left, content_top, chart_w, content_height


# ---------------------------------------------------------------------------
# Category chart
# ---------------------------------------------------------------------------

def _add_category_chart(slide, chart_el, left, top, width, height):
    xl_type = CATEGORY_CHART_MAP.get(chart_el.chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
    chart_data = CategoryChartData()
    chart_data.categories = chart_el.categories
    for s in chart_el.series:
        chart_data.add_series(s.name, s.values)

    shape = slide.shapes.add_chart(xl_type, left, top, width, height, chart_data)
    chart = shape.chart
    _set_chart_title(chart, chart_el.title)
    chart.has_legend = chart_el.legend or len(chart_el.series) > 1
    plot = chart.plots[0]

    if chart_el.labels:
        _apply_data_labels(plot, chart_el.chart_type, chart_el.number_format)

    _apply_series_colors(plot, chart_el.series)
    return shape


# ---------------------------------------------------------------------------
# Scatter chart
# ---------------------------------------------------------------------------

def _add_scatter_chart(slide, chart_el, left, top, width, height):
    chart_data = XyChartData()
    for s in chart_el.series:
        series_data = chart_data.add_series(s.name)
        xs = s.x_values or []
        ys = s.y_values or []
        for x, y in zip(xs, ys):
            series_data.add_data_point(x, y)

    shape = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS,
        left, top, width, height, chart_data,
    )
    chart = shape.chart
    _set_chart_title(chart, chart_el.title)
    chart.has_legend = chart_el.legend or len(chart_el.series) > 1
    plot = chart.plots[0]
    _apply_series_colors(plot, chart_el.series)
    return shape


# ---------------------------------------------------------------------------
# Bubble chart
# ---------------------------------------------------------------------------

def _add_bubble_chart(slide, chart_el, left, top, width, height):
    chart_data = BubbleChartData()
    for s in chart_el.series:
        series_data = chart_data.add_series(s.name)
        xs = s.x_values or []
        ys = s.y_values or []
        szs = s.sizes or [1] * len(xs)
        for x, y, sz in zip(xs, ys, szs):
            series_data.add_data_point(x, y, sz)

    shape = slide.shapes.add_chart(
        XL_CHART_TYPE.BUBBLE,
        left, top, width, height, chart_data,
    )
    chart = shape.chart
    _set_chart_title(chart, chart_el.title)
    chart.has_legend = chart_el.legend or len(chart_el.series) > 1
    plot = chart.plots[0]
    _apply_series_colors(plot, chart_el.series)
    return shape


# ---------------------------------------------------------------------------
# Waterfall chart (stacked column illusion)
# ---------------------------------------------------------------------------

def _waterfall_series(values: list[float], totals: list[int]):
    """
    Decompose waterfall values into (base, gain, loss) series.

    For normal bars: base = running cumulative, gain/loss = positive/negative delta.
    For total bars: base = 0, gain = total (shown from zero).
    Returns three lists of same length as values.
    """
    n = len(values)
    base_vals  = [0.0] * n
    gain_vals  = [0.0] * n
    loss_vals  = [0.0] * n

    running = 0.0
    for i, v in enumerate(values):
        if i in totals:
            gain_vals[i] = running + v if v >= 0 else 0
            loss_vals[i] = 0 if v >= 0 else -(running + v)
            running = running + v
        else:
            if v >= 0:
                base_vals[i] = running
                gain_vals[i] = v
                running += v
            else:
                base_vals[i] = running + v
                loss_vals[i] = -v
                running += v

    return base_vals, gain_vals, loss_vals


def _make_transparent(series_xml):
    """Set a series fill to no-fill (transparent) via XML."""
    # Remove existing spPr if any, add a new one with noFill
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    c_ns = "http://schemas.openxmlformats.org/drawingml/2006/chart"
    existing = series_xml.find(qn("c:spPr"))
    if existing is not None:
        series_xml.remove(existing)
    spPr = etree.SubElement(series_xml, qn("c:spPr"))
    noFill = etree.SubElement(spPr, qn("a:noFill"))


def _set_series_fill_color(series_xml, rgb: RGBColor):
    existing = series_xml.find(qn("c:spPr"))
    if existing is not None:
        series_xml.remove(existing)
    spPr = etree.SubElement(series_xml, qn("c:spPr"))
    solidFill = etree.SubElement(spPr, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")


def _add_waterfall_chart(slide, chart_el, left, top, width, height):
    values = chart_el.series[0].values if chart_el.series else []
    if not values:
        return None
    totals = chart_el.totals
    colors = chart_el.colors or {}

    gain_color  = _parse_color(colors.get("gain",  "")) or _WF_GAIN_COLOR
    loss_color  = _parse_color(colors.get("loss",  "")) or _WF_LOSS_COLOR
    total_color = _parse_color(colors.get("total", "")) or _WF_TOTAL_COLOR

    base_vals, gain_vals, loss_vals = _waterfall_series(values, totals)

    chart_data = CategoryChartData()
    chart_data.categories = chart_el.categories
    chart_data.add_series("_base", base_vals)
    chart_data.add_series("增量", gain_vals)
    chart_data.add_series("减量", loss_vals)

    shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED,
        left, top, width, height, chart_data,
    )
    chart = shape.chart
    _set_chart_title(chart, chart_el.title)
    chart.has_legend = False

    # Apply colors via XML: base = transparent, gain = green, loss = red
    plot_xml = chart._element.find(
        ".//" + qn("c:barChart")
    )
    if plot_xml is not None:
        ser_list = plot_xml.findall(qn("c:ser"))
        if len(ser_list) >= 1:
            _make_transparent(ser_list[0])
        if len(ser_list) >= 2:
            _set_series_fill_color(ser_list[1], gain_color)
            # override total bars to total_color
            for idx in totals:
                _add_individual_point_color(ser_list[1], idx, total_color)
        if len(ser_list) >= 3:
            _set_series_fill_color(ser_list[2], loss_color)

    # Data labels on gain/loss (not base)
    if chart_el.labels:
        plot = chart.plots[0]
        _apply_data_labels(plot, "column", chart_el.number_format)

    return shape


def _add_individual_point_color(ser_xml, pt_idx: int, rgb: RGBColor):
    """Override fill color for a single data point in a series via OOXML."""
    dPt = etree.SubElement(ser_xml, qn("c:dPt"))
    idx_el = etree.SubElement(dPt, qn("c:idx"))
    idx_el.set("val", str(pt_idx))
    spPr = etree.SubElement(dPt, qn("c:spPr"))
    solidFill = etree.SubElement(spPr, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")


# ---------------------------------------------------------------------------
# Combo chart (bar + line, with optional secondary axis)
# ---------------------------------------------------------------------------

def _build_ser_xml(idx: int, name: str, categories: list, values: list[float],
                   ax_id_cat: int, ax_id_val: int) -> etree._Element:
    """Build a minimal <c:ser> element for a category series."""
    C = "http://schemas.openxmlformats.org/drawingml/2006/chart"
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"

    ser = etree.Element(qn("c:ser"))
    _sub(ser, "c:idx", val=str(idx))
    _sub(ser, "c:order", val=str(idx))

    # Series name
    tx = etree.SubElement(ser, qn("c:tx"))
    strRef = etree.SubElement(tx, qn("c:strRef"))
    _sub(strRef, "c:f", text=f"Sheet1!${_col_letter(idx+1)}$1")
    strCache = etree.SubElement(strRef, qn("c:strCache"))
    _sub(strCache, "c:ptCount", val="1")
    pt = etree.SubElement(strCache, qn("c:pt"))
    pt.set("idx", "0")
    _sub(pt, "c:v", text=str(name))

    # Categories
    cat = etree.SubElement(ser, qn("c:cat"))
    catStrRef = etree.SubElement(cat, qn("c:strRef"))
    _sub(catStrRef, "c:f", text=f"Sheet1!$A$2:$A${len(categories)+1}")
    catCache = etree.SubElement(catStrRef, qn("c:strCache"))
    _sub(catCache, "c:ptCount", val=str(len(categories)))
    for ci, cat_name in enumerate(categories):
        cpt = etree.SubElement(catCache, qn("c:pt"))
        cpt.set("idx", str(ci))
        _sub(cpt, "c:v", text=str(cat_name))

    # Values
    val = etree.SubElement(ser, qn("c:val"))
    numRef = etree.SubElement(val, qn("c:numRef"))
    _sub(numRef, "c:f", text=f"Sheet1!${_col_letter(idx+1)}$2:${_col_letter(idx+1)}${len(values)+1}")
    numCache = etree.SubElement(numRef, qn("c:numCache"))
    _sub(numCache, "c:formatCode", text="General")
    _sub(numCache, "c:ptCount", val=str(len(values)))
    for vi, v in enumerate(values):
        vpt = etree.SubElement(numCache, qn("c:pt"))
        vpt.set("idx", str(vi))
        _sub(vpt, "c:v", text=str(v))

    return ser


def _sub(parent, tag: str, val: str = None, text: str = None) -> etree._Element:
    el = etree.SubElement(parent, qn(tag))
    if val is not None:
        el.set("val", val)
    if text is not None:
        el.text = text
    return el


def _col_letter(n: int) -> str:
    """1 → A, 2 → B, ... 26 → Z, 27 → AA"""
    result = ""
    while n > 0:
        n, r = divmod(n - 1, 26)
        result = chr(65 + r) + result
    return result


def _add_combo_chart(slide, chart_el, left, top, width, height):
    """
    Build a combo chart by:
    1. Creating a COLUMN chart with all bar series
    2. Doing XML surgery: add a <c:lineChart> for line series,
       optionally add a secondary valAx for series marked axis=secondary
    """
    categories = chart_el.categories
    bar_series  = [s for s in chart_el.series if s.subtype != "line"]
    line_series = [s for s in chart_el.series if s.subtype == "line"]

    # Start with bar series only
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for s in bar_series:
        chart_data.add_series(s.name, s.values)

    shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        left, top, width, height, chart_data,
    )
    chart = shape.chart
    _set_chart_title(chart, chart_el.title)
    chart.has_legend = True

    # Apply colors to bar series
    plot_xml = chart._element.find(".//" + qn("c:barChart"))
    if plot_xml is not None:
        _apply_series_colors_xml(plot_xml, bar_series)

    if not line_series:
        return shape

    # Get existing axis IDs
    plot_area = chart._element.find(".//" + qn("c:plotArea"))
    bar_chart = plot_area.find(qn("c:barChart"))
    axId_elems = bar_chart.findall(qn("c:axId"))
    pri_cat_ax_id = int(axId_elems[0].get("val"))
    pri_val_ax_id = int(axId_elems[1].get("val"))

    # Determine if we need a secondary axis
    needs_secondary = any(s.axis == "secondary" for s in line_series)
    sec_val_ax_id = pri_val_ax_id + 1000 if needs_secondary else pri_val_ax_id

    # Build <c:lineChart>
    lineChart = etree.Element(qn("c:lineChart"))
    _sub(lineChart, "c:grouping", val="standard")

    total_series_count = len(bar_series)
    for i, s in enumerate(line_series):
        ser_idx = total_series_count + i
        use_ax = sec_val_ax_id if s.axis == "secondary" else pri_val_ax_id
        ser_xml = _build_ser_xml(ser_idx, s.name, categories, s.values,
                                 pri_cat_ax_id, use_ax)
        # Add line marker
        marker = etree.SubElement(ser_xml, qn("c:marker"))
        _sub(marker, "c:symbol", val="circle")
        lineChart.append(ser_xml)

        if s.color:
            _set_series_fill_color(ser_xml, _parse_color(s.color))

    ax_to_use = sec_val_ax_id if needs_secondary else pri_val_ax_id
    _sub(lineChart, "c:axId", val=str(pri_cat_ax_id))
    _sub(lineChart, "c:axId", val=str(ax_to_use))

    # Insert lineChart into plotArea (before axes)
    cat_ax = plot_area.find(qn("c:catAx"))
    if cat_ax is not None:
        cat_ax.addprevious(lineChart)
    else:
        plot_area.append(lineChart)

    # Add secondary valAx if needed
    if needs_secondary:
        # Copy the primary valAx and modify for secondary
        pri_val_ax = plot_area.find(qn("c:valAx"))
        sec_val_ax = copy.deepcopy(pri_val_ax)
        # Update axId
        sec_val_ax.find(qn("c:axId")).set("val", str(sec_val_ax_id))
        # Place on right
        axPos = sec_val_ax.find(qn("c:axPos"))
        if axPos is not None:
            axPos.set("val", "r")
        # Hide gridlines on secondary
        for gl in sec_val_ax.findall(qn("c:majorGridlines")):
            sec_val_ax.remove(gl)
        # Update crossAx to primary cat
        crossAx = sec_val_ax.find(qn("c:crossAx"))
        if crossAx is not None:
            crossAx.set("val", str(pri_cat_ax_id))
        plot_area.append(sec_val_ax)

    return shape


def _apply_series_colors_xml(chart_type_xml, series_list):
    """Apply colors to series in a barChart/lineChart XML element."""
    ser_elems = chart_type_xml.findall(qn("c:ser"))
    for i, sd in enumerate(series_list):
        if sd.color and i < len(ser_elems):
            color = _parse_color(sd.color)
            if color:
                _set_series_fill_color(ser_elems[i], color)


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def add_chart_to_slide(slide, chart_el, left, top, width, height):
    """Dispatch to the appropriate chart builder based on chart_type."""
    ct = chart_el.chart_type

    if ct == "waterfall":
        return _add_waterfall_chart(slide, chart_el, left, top, width, height)

    if ct == "combo":
        return _add_combo_chart(slide, chart_el, left, top, width, height)

    if ct == "scatter":
        return _add_scatter_chart(slide, chart_el, left, top, width, height)

    if ct == "bubble":
        return _add_bubble_chart(slide, chart_el, left, top, width, height)

    # All remaining types use CategoryChartData
    if ct in CATEGORY_CHART_MAP:
        return _add_category_chart(slide, chart_el, left, top, width, height)

    # Unknown type → fallback to clustered column
    print(f"  Warning: unknown chart type '{ct}', falling back to column")
    chart_el.chart_type = "column"
    return _add_category_chart(slide, chart_el, left, top, width, height)
