#!/usr/bin/env python3
"""Extract chart data from an existing PPTX as chart YAML blocks.

Usage:
    python3 scripts/extract_charts.py input.pptx [-o output.md]

Produces markdown-ready ```chart blocks that can be pasted into slide markdown.
"""

from __future__ import annotations

import argparse
import os
import sys

from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE


# Reverse map from python-pptx chart type enum to our simple names
_CHART_TYPE_REVERSE = {
    XL_CHART_TYPE.COLUMN_CLUSTERED: "column",
    XL_CHART_TYPE.COLUMN_STACKED: "column",
    XL_CHART_TYPE.BAR_CLUSTERED: "bar",
    XL_CHART_TYPE.BAR_STACKED: "bar",
    XL_CHART_TYPE.LINE: "line",
    XL_CHART_TYPE.LINE_MARKERS: "line",
    XL_CHART_TYPE.PIE: "pie",
}


def extract_charts(pptx_path: str) -> list[dict]:
    """Extract all charts from a PPTX, returning structured data."""
    prs = Presentation(pptx_path)
    results = []

    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not (hasattr(shape, 'has_chart') and shape.has_chart):
                continue
            chart = shape.chart
            simple_type = _CHART_TYPE_REVERSE.get(chart.chart_type, "column")

            plot = chart.plots[0]
            categories = [str(c) for c in plot.categories]
            series_list = []
            for series in plot.series:
                series_list.append({
                    'name': str(series.tx if hasattr(series, 'tx') else f'Series {len(series_list)+1}'),
                    'values': list(series.values),
                })

            title = ""
            if chart.has_title:
                title = chart.chart_title.text_frame.text

            results.append({
                'slide': si + 1,
                'shape_name': shape.name,
                'type': simple_type,
                'title': title,
                'categories': categories,
                'series': series_list,
            })

    return results


def charts_to_markdown(charts: list[dict]) -> str:
    """Convert extracted chart data to markdown chart blocks."""
    lines = []
    current_slide = None

    for c in charts:
        if c['slide'] != current_slide:
            current_slide = c['slide']
            lines.append(f"\n<!-- Slide {current_slide}: {c['shape_name']} -->")

        lines.append("")
        lines.append("```chart")
        lines.append(f"type: {c['type']}")
        if c['title']:
            lines.append(f"title: \"{c['title']}\"")

        cats_str = ', '.join(f'"{cat}"' for cat in c['categories'])
        lines.append(f"categories: [{cats_str}]")
        lines.append("series:")
        for s in c['series']:
            lines.append(f"  - name: \"{s['name']}\"")
            vals_str = ', '.join(str(v) for v in s['values'])
            lines.append(f"    values: [{vals_str}]")
        lines.append("```")

    return '\n'.join(lines)


def main():
    parser = argparse.ArgumentParser(description='Extract charts from PPTX')
    parser.add_argument('input', help='Input PPTX file')
    parser.add_argument('-o', '--output', default=None,
                        help='Output markdown file (default: stdout)')
    args = parser.parse_args()

    if not os.path.exists(args.input):
        print(f"Error: {args.input} not found", file=sys.stderr)
        sys.exit(1)

    charts = extract_charts(args.input)
    md = charts_to_markdown(charts)

    if args.output:
        with open(args.output, 'w', encoding='utf-8') as f:
            f.write(md)
        print(f"Extracted {len(charts)} charts to {args.output}")
    else:
        print(md)
        print(f"\n# Extracted {len(charts)} charts", file=sys.stderr)


if __name__ == '__main__':
    main()
