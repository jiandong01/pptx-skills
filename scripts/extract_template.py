#!/usr/bin/env python3
"""Extract a clean template from an existing pptx file.

Usage:
    python3 scripts/extract_template.py input.pptx [-o output_dir]

Outputs:
    output_dir/template.pptx  - Empty template with masters/layouts/theme
    output_dir/template.md    - Markdown skeleton with layout documentation
"""

from __future__ import annotations

import argparse
import copy
import os
import sys

from pptx import Presentation

from slide_utils import get_layout_info, get_used_layouts, wps_fixup


def strip_slides(src_path: str, dst_path: str):
    """Load a pptx, remove all slides, save as empty template."""
    prs = Presentation(src_path)

    # Collect layout usage info before stripping
    used_layouts = get_used_layouts(prs)
    all_layouts = get_layout_info(prs)

    # Remove all slides via XML manipulation
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        if rId:
            prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)

    prs.save(dst_path)

    # WPS compatibility fixup
    wps_fixup(dst_path)

    return all_layouts, used_layouts


def generate_template_md(output_path: str, all_layouts: list[dict], used_layouts: list[dict]):
    """Generate a markdown skeleton documenting available layouts."""
    lines = [
        '---',
        'title: "Presentation Title"',
        'subtitle: "Subtitle"',
        'template: "templates/template.pptx"',
        '---',
        '',
        '<!--',
        '  Available Layouts:',
    ]

    # Group by master
    masters = {}
    for layout in all_layouts:
        mi = layout['master']
        if mi not in masters:
            masters[mi] = []
        masters[mi].append(layout)

    for mi in sorted(masters.keys()):
        lines.append(f'')
        lines.append(f'  Master {mi}:')
        for layout in masters[mi]:
            ph_info = ', '.join(
                f"idx={p['idx']}({p['name']})" for p in layout['placeholders']
            )
            lines.append(f"    - \"{layout['name']}\": [{ph_info}]")

    lines.append('')
    lines.append('  Used in source presentation:')
    for u in used_layouts:
        lines.append(f"    - Master {u['master']}: \"{u['name']}\"")

    lines.append('-->')
    lines.append('')

    # Generate example slides for commonly used layouts
    lines.append('## Title Page')
    lines.append('<!-- layout: "Title Slide" -->')
    lines.append('')
    lines.append('---')
    lines.append('')
    lines.append('## Content Slide')
    lines.append('<!-- layout: "Title and Content" -->')
    lines.append('')
    lines.append('Slide body text with **bold** and *italic* support.')
    lines.append('')
    lines.append('- Bullet point 1')
    lines.append('- Bullet point 2')
    lines.append('  - Sub-bullet')
    lines.append('')
    lines.append('---')
    lines.append('')
    lines.append('## Content with Caption Slide')
    lines.append('<!-- layout: "Content with Caption" -->')
    lines.append('')
    lines.append('> Caption or description text (goes to the caption area)')
    lines.append('')
    lines.append('| Header 1 | Header 2 |')
    lines.append('|----------|----------|')
    lines.append('| Cell 1   | Cell 2   |')
    lines.append('')

    with open(output_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(lines) + '\n')


def main():
    parser = argparse.ArgumentParser(description='Extract template from pptx')
    parser.add_argument('input', help='Input pptx file')
    parser.add_argument('-o', '--output-dir', default='templates',
                        help='Output directory (default: templates/)')
    args = parser.parse_args()

    if not os.path.exists(args.input):
        print(f"Error: {args.input} not found", file=sys.stderr)
        sys.exit(1)

    os.makedirs(args.output_dir, exist_ok=True)

    tpl_path = os.path.join(args.output_dir, 'template.pptx')
    md_path = os.path.join(args.output_dir, 'template.md')

    print(f"==> Extracting template from {args.input}...")
    all_layouts, used_layouts = strip_slides(args.input, tpl_path)
    print(f"  Template saved: {tpl_path}")
    print(f"  Masters: {len(set(l['master'] for l in all_layouts))}")
    print(f"  Total layouts: {len(all_layouts)}")
    print(f"  Used layouts: {len(used_layouts)}")

    print(f"==> Generating markdown skeleton...")
    generate_template_md(md_path, all_layouts, used_layouts)
    print(f"  Skeleton saved: {md_path}")

    print(f"\n==> Done. Layouts found:")
    for layout in used_layouts:
        ph_str = ', '.join(f"idx={p['idx']}" for p in layout['placeholders'])
        print(f"  Master {layout['master']}: \"{layout['name']}\" [{ph_str}]")


if __name__ == '__main__':
    main()
