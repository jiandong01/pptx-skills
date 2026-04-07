#!/usr/bin/env python3
"""
列出 PowerPoint 模板的所有布局

用法:
    python list_layouts.py <template.pptx>

示例:
    python list_layouts.py ../template/PKU/近期系统安全相关工作.pptx
"""

import sys
from pathlib import Path
from pptx import Presentation

# 添加当前目录到 Python 路径
sys.path.insert(0, str(Path(__file__).parent))

from layout_standards import find_standard_by_keywords, STANDARD_LAYOUTS


def list_layouts(template_path: str):
    """列出模板的所有布局及其推荐用法"""
    try:
        prs = Presentation(template_path)
    except Exception as e:
        print(f"错误: 无法打开模板文件 '{template_path}'")
        print(f"原因: {e}")
        return

    print(f"\n{'='*70}")
    print(f"模板布局列表: {Path(template_path).name}")
    print(f"{'='*70}")
    print(f"幻灯片尺寸: {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")
    print(f"总布局数: {len(prs.slide_master.slide_layouts)}")
    print(f"{'='*70}\n")

    for idx, layout in enumerate(prs.slide_master.slide_layouts):
        # 基本信息
        print(f"[{idx}] {layout.name}")

        # 占位符信息
        if layout.placeholders:
            ph_info = []
            for ph in layout.placeholders:
                ph_type = ph.placeholder_format.type
                ph_name = ph.name
                # 简化占位符类型显示
                type_map = {
                    1: "标题",
                    2: "正文",
                    3: "中心标题",
                    4: "副标题",
                    7: "日期",
                    8: "页脚",
                    10: "页码",
                    11: "图片",
                    12: "表格",
                }
                type_str = type_map.get(ph_type, f"类型{ph_type}")
                ph_info.append(f"{ph_name}({type_str})")

            print(f"    占位符: {', '.join(ph_info[:4])}", end="")
            if len(ph_info) > 4:
                print(f" ... 共{len(ph_info)}个")
            else:
                print()
        else:
            print(f"    占位符: 无")

        # 推荐的标准布局别名
        std_name = find_standard_by_keywords(layout.name)
        if std_name:
            std = STANDARD_LAYOUTS[std_name]
            print(f"    推荐别名: {std_name}")
            print(f"    其他别名: {', '.join(std.aliases)}")
            print(f"    用途: {std.description}")
        else:
            print(f"    推荐别名: (无标准映射)")

        print()

    # 使用说明
    print(f"{'='*70}")
    print("使用方法:")
    print()
    print("1. 在 Markdown 中使用布局索引:")
    print("   ## 我的幻灯片")
    print("   <!-- layout: 7 -->")
    print()
    print("2. 在 Markdown 中使用标准别名:")
    print("   ## 我的幻灯片")
    print("   <!-- layout: standard -->")
    print()
    print("3. 在 Markdown 中使用模板布局名称:")
    print("   ## 我的幻灯片")
    print("   <!-- layout: Title and Content -->")
    print()
    print("标准布局别名列表:")
    print("  结构性: cover, toc, section, summary")
    print("  内容性: standard, two-column, image, chart, table, mixed, title-only")
    print(f"{'='*70}\n")


def main():
    if len(sys.argv) < 2:
        print("用法: python list_layouts.py <template.pptx>")
        print()
        print("示例:")
        print("  python list_layouts.py ../template/PKU/近期系统安全相关工作.pptx")
        print("  python list_layouts.py ../template/default.pptx")
        sys.exit(1)

    template_path = sys.argv[1]
    if not Path(template_path).exists():
        print(f"错误: 文件不存在 '{template_path}'")
        sys.exit(1)

    list_layouts(template_path)


if __name__ == '__main__':
    main()
