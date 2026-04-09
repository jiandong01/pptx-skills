"""Shared utilities for pptx slide generation from markdown.

Includes: markdown parser, layout finder, text formatting, WPS fixup, mermaid rendering.
"""

from __future__ import annotations

import os
import re
import subprocess
import tempfile
import zipfile
from dataclasses import dataclass, field
from pathlib import Path

import yaml
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu


# ---------------------------------------------------------------------------
# Data structures
# ---------------------------------------------------------------------------

@dataclass
class Run:
    text: str
    bold: bool = False
    italic: bool = False


@dataclass
class Paragraph:
    runs: list[Run] = field(default_factory=list)
    level: int = 0


@dataclass
class TextElement:
    paragraphs: list[Paragraph] = field(default_factory=list)


@dataclass
class BlockquoteElement:
    paragraphs: list[Paragraph] = field(default_factory=list)


@dataclass
class TableElement:
    headers: list[str] = field(default_factory=list)
    rows: list[list[str]] = field(default_factory=list)


@dataclass
class ImageElement:
    path: str
    alt: str = ""


@dataclass
class MermaidElement:
    code: str
    index: int = 0


@dataclass
class ChartSeries:
    name: str = ""
    values: list[float] = field(default_factory=list)
    color: str = ""          # hex color e.g. "#C00000", empty = auto
    # scatter/bubble
    x_values: list[float] = field(default_factory=list)
    y_values: list[float] = field(default_factory=list)
    sizes: list[float] = field(default_factory=list)   # bubble only
    # combo
    subtype: str = ""        # "bar" | "line" — for combo charts
    axis: str = "primary"    # "primary" | "secondary" — for combo line series
    number_format: str = ""  # per-series format override


@dataclass
class ChartElement:
    chart_type: str = "column"
    title: str = ""
    categories: list[str] = field(default_factory=list)
    series: list[ChartSeries] = field(default_factory=list)
    position: str = "center"   # left | right | center
    width: str = "60%"
    labels: bool = True
    legend: bool = False
    number_format: str = ""    # global format
    # waterfall
    totals: list[int] = field(default_factory=list)   # index positions of total bars
    colors: dict = field(default_factory=dict)         # {gain, loss, total} overrides


@dataclass
class SlideData:
    title: str = ""
    layout_hint: str | None = None
    body_elements: list = field(default_factory=list)


@dataclass
class PresentationData:
    title: str = ""
    subtitle: str = ""
    template: str = "templates/template.pptx"
    output: str = ""
    slides: list[SlideData] = field(default_factory=list)


# ---------------------------------------------------------------------------
# Markdown parser
# ---------------------------------------------------------------------------

def parse_inline(text: str) -> list[Run]:
    """Parse inline markdown formatting (**bold**, *italic*) into Runs."""
    runs = []
    # Pattern: **bold**, *italic*, or plain text
    pattern = re.compile(r'(\*\*(.+?)\*\*|\*(.+?)\*|([^*]+))')
    for m in pattern.finditer(text):
        if m.group(2) is not None:
            runs.append(Run(text=m.group(2), bold=True))
        elif m.group(3) is not None:
            runs.append(Run(text=m.group(3), italic=True))
        elif m.group(4) is not None:
            runs.append(Run(text=m.group(4)))
    if not runs and text:
        runs.append(Run(text=text))
    return runs


def _bullet_level(line: str) -> tuple[int, str]:
    """Return (indent_level, text) for a bullet line like '- text' or '  - text'."""
    stripped = line.lstrip()
    indent = len(line) - len(stripped)
    level = indent // 2
    text = stripped.lstrip("- ").strip()
    return level, text


def _is_table_separator(line: str) -> bool:
    """Check if line is a markdown table separator like |---|---|."""
    return bool(re.match(r'^\s*\|[\s\-:|]+\|\s*$', line))


def _parse_table_row(line: str) -> list[str]:
    """Parse a markdown table row into cell texts."""
    cells = line.strip().strip('|').split('|')
    return [c.strip() for c in cells]


def parse_markdown(md_path: str) -> PresentationData:
    """Parse a slide markdown file into PresentationData."""
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()

    pdata = PresentationData()
    pdata.output = str(Path(md_path).with_suffix('.pptx'))

    # Extract YAML front matter
    fm_match = re.match(r'^---\s*\n(.*?)\n---\s*\n', content, re.DOTALL)
    body = content
    if fm_match:
        fm = yaml.safe_load(fm_match.group(1))
        if fm:
            pdata.title = fm.get('title', '')
            pdata.subtitle = fm.get('subtitle', '').replace('\\n', '\n')
            pdata.template = fm.get('template', pdata.template)
            pdata.output = fm.get('output', pdata.output)
        body = content[fm_match.end():]

    # Split into slides by horizontal rule
    raw_slides = re.split(r'\n---\s*\n', body)

    mermaid_idx = 0
    for raw in raw_slides:
        raw = raw.strip()
        if not raw:
            continue

        slide = SlideData()
        lines = raw.split('\n')
        body_lines = []

        # Extract title (## heading)
        for i, line in enumerate(lines):
            if line.startswith('## '):
                slide.title = line[3:].strip()
                body_lines = lines[i + 1:]
                break
        else:
            body_lines = lines

        # Extract layout hint <!-- layout: "..." -->
        new_body = []
        for line in body_lines:
            lh_match = re.match(r'^\s*<!--\s*layout:\s*["\']?(.+?)["\']?\s*-->\s*$', line)
            if lh_match:
                slide.layout_hint = lh_match.group(1).strip()
            else:
                new_body.append(line)
        body_lines = new_body

        # Parse body into elements
        slide.body_elements = _parse_body(body_lines, mermaid_idx)
        # Update mermaid counter
        for el in slide.body_elements:
            if isinstance(el, MermaidElement):
                mermaid_idx = el.index + 1

        pdata.slides.append(slide)

    return pdata


def _parse_body(lines: list[str], mermaid_idx: int) -> list:
    """Parse body lines into a list of content elements."""
    elements = []
    i = 0

    while i < len(lines):
        line = lines[i]

        # Skip blank lines
        if not line.strip():
            i += 1
            continue

        # Skip HTML comments (non-layout)
        if re.match(r'^\s*<!--.*-->\s*$', line):
            i += 1
            continue

        # Mermaid code block
        if line.strip().startswith('```mermaid'):
            code_lines = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith('```'):
                code_lines.append(lines[i])
                i += 1
            i += 1  # skip closing ```
            elements.append(MermaidElement(code='\n'.join(code_lines), index=mermaid_idx))
            mermaid_idx += 1
            continue

        # Chart code block
        if line.strip().startswith('```chart'):
            code_lines = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith('```'):
                code_lines.append(lines[i])
                i += 1
            i += 1  # skip closing ```
            chart_yaml = yaml.safe_load('\n'.join(code_lines))
            if chart_yaml:
                chart_el = ChartElement(
                    chart_type=chart_yaml.get('type', 'column'),
                    title=chart_yaml.get('title', ''),
                    categories=chart_yaml.get('categories', []),
                    position=chart_yaml.get('position', 'center'),
                    width=chart_yaml.get('width', '60%'),
                    labels=chart_yaml.get('labels', True),
                    legend=chart_yaml.get('legend', False),
                    number_format=chart_yaml.get('number_format', ''),
                    totals=chart_yaml.get('totals', []),
                    colors=chart_yaml.get('colors', {}),
                )
                for s in chart_yaml.get('series', []):
                    chart_el.series.append(ChartSeries(
                        name=s.get('name', ''),
                        values=s.get('values', []),
                        color=s.get('color', ''),
                        x_values=s.get('x_values', []),
                        y_values=s.get('y_values', []),
                        sizes=s.get('sizes', []),
                        subtype=s.get('subtype', ''),
                        axis=s.get('axis', 'primary'),
                        number_format=s.get('number_format', ''),
                    ))
                elements.append(chart_el)
            continue

        # Other code blocks - skip
        if line.strip().startswith('```'):
            i += 1
            while i < len(lines) and not lines[i].strip().startswith('```'):
                i += 1
            i += 1
            continue

        # Blockquote — match `> text` and bare `>` (empty blockquote line)
        if line.strip().startswith('>'):
            bq_paras = []
            while i < len(lines) and lines[i].strip().startswith('>'):
                stripped = lines[i].strip()
                text = stripped[2:] if stripped.startswith('> ') else stripped[1:]
                bq_paras.append(Paragraph(runs=parse_inline(text)))
                i += 1
            elements.append(BlockquoteElement(paragraphs=bq_paras))
            continue

        # Table
        if '|' in line and i + 1 < len(lines) and _is_table_separator(lines[i + 1]):
            headers = _parse_table_row(line)
            i += 2  # skip header + separator
            rows = []
            while i < len(lines) and '|' in lines[i] and lines[i].strip():
                rows.append(_parse_table_row(lines[i]))
                i += 1
            elements.append(TableElement(headers=headers, rows=rows))
            continue

        # Image
        img_match = re.match(r'^\s*!\[([^\]]*)\]\(([^)]+)\)\s*$', line)
        if img_match:
            elements.append(ImageElement(alt=img_match.group(1), path=img_match.group(2)))
            i += 1
            continue

        # Bullet list
        if re.match(r'^\s*[-*]\s', line):
            paras = []
            while i < len(lines) and re.match(r'^\s*[-*]\s', lines[i]):
                level, text = _bullet_level(lines[i])
                paras.append(Paragraph(runs=parse_inline(text), level=level))
                i += 1
            elements.append(TextElement(paragraphs=paras))
            continue

        # Ordered list
        if re.match(r'^\s*\d+[.)]\s', line):
            paras = []
            while i < len(lines) and re.match(r'^\s*\d+[.)]\s', lines[i]):
                stripped = lines[i].strip()
                text = re.sub(r'^\d+[.)]\s+', '', stripped)
                level = (len(lines[i]) - len(lines[i].lstrip())) // 2
                paras.append(Paragraph(runs=parse_inline(text), level=level))
                i += 1
            elements.append(TextElement(paragraphs=paras))
            continue

        # Plain text paragraph (collect consecutive non-empty, non-special lines)
        para_lines = []
        while (i < len(lines) and lines[i].strip()
               and not lines[i].strip().startswith(('#', '>', '|', '-', '```', '!'))
               and not re.match(r'^\s*\d+[.)]\s', lines[i])):
            para_lines.append(lines[i].strip())
            i += 1
        if para_lines:
            text = ' '.join(para_lines)
            elements.append(TextElement(paragraphs=[Paragraph(runs=parse_inline(text))]))
        else:
            i += 1  # safety: ensure forward progress if no branch consumed the line

    return elements


# ---------------------------------------------------------------------------
# Layout finder
# ---------------------------------------------------------------------------

from layout_standards import (
    STANDARD_LAYOUTS,
    get_layout_standard,
    find_standard_by_keywords
)


def find_layout_flexible(prs: Presentation, hint: str) -> tuple:
    """灵活查找布局：支持索引、标准名称、别名、模板名称

    Returns:
        (layout, standard_name) - 布局对象和对应的标准名称（如果有）
    """
    # 1. 尝试按索引查找
    try:
        idx = int(hint)
        layouts = prs.slide_master.slide_layouts
        if 0 <= idx < len(layouts):
            # 尝试推断标准名称
            layout = layouts[idx]
            std_name = find_standard_by_keywords(layout.name)
            return (layout, std_name)
    except (ValueError, AttributeError):
        pass

    # 2. 检查是否为标准布局名称或别名
    layout_std = get_layout_standard(hint)
    if layout_std:
        # 在模板中查找匹配的布局
        layout = _find_by_keywords(prs, layout_std.keywords)
        if layout:
            return (layout, layout_std.name)

        # 找不到，尝试回退策略
        fallback = _get_fallback_layout(prs, layout_std.name)
        if fallback:
            return fallback

    # 3. 按模板布局名称精确匹配
    for layout in prs.slide_master.slide_layouts:
        if layout.name.lower() == hint.lower():
            std_name = find_standard_by_keywords(layout.name)
            return (layout, std_name)

    # 4. 按模板布局名称模糊匹配
    hint_lower = hint.lower()
    for layout in prs.slide_master.slide_layouts:
        if hint_lower in layout.name.lower():
            std_name = find_standard_by_keywords(layout.name)
            return (layout, std_name)

    # 5. 找不到，返回默认布局
    print(f"Warning: 找不到布局 '{hint}'，使用默认布局")
    default_layout = _get_default_layout(prs)
    return (default_layout, 'standard')


def _find_by_keywords(prs: Presentation, keywords: list[str]):
    """根据关键词列表查找布局"""
    for keyword in keywords:
        keyword_lower = keyword.lower()
        for layout in prs.slide_master.slide_layouts:
            if keyword_lower in layout.name.lower():
                return layout
    return None


def _get_fallback_layout(prs: Presentation, std_name: str) -> tuple:
    """获取回退布局"""
    fallback_map = {
        'cover': ['section', 'title-only', 'standard'],
        'toc': ['standard', 'title-only'],
        'section': ['cover', 'title-only', 'standard'],
        'summary': ['section', 'cover', 'standard'],
        'image': ['title-only', 'standard'],
        'chart': ['title-only', 'standard'],
        'table': ['standard', 'title-only'],
        'mixed': ['two-column', 'standard'],        'two-column': ['standard', 'title-only'],
    }

    fallbacks = fallback_map.get(std_name, ['standard'])
    for fallback_name in fallbacks:
        fallback_std = get_layout_standard(fallback_name)
        if fallback_std:
            layout = _find_by_keywords(prs, fallback_std.keywords)
            if layout:
                print(f"  使用回退布局: {fallback_name} ({layout.name})")
                return (layout, std_name)

    # 最后的保底
    default_layout = _get_default_layout(prs)
    return (default_layout, 'standard')


def _get_default_layout(prs: Presentation):
    """Return the best fallback layout: prefer 'Title and Content', else first with a title ph."""
    for layout in prs.slide_master.slide_layouts:
        name_lower = layout.name.lower()
        if 'title' in name_lower and 'content' in name_lower:
            return layout
        if any(ph.placeholder_format.type == 1 for ph in layout.placeholders):
            return layout
    return prs.slide_master.slide_layouts[0]


def find_layout(prs: Presentation, name: str):
    """Thin wrapper around find_layout_flexible for callers that only need the layout."""
    layout, _ = find_layout_flexible(prs, name)
    return layout


def select_layout(slide: SlideData, prs: Presentation) -> tuple:
    """Select the appropriate layout for a slide.

    Returns:
        (layout, std_name) — layout 对象 + 标准布局名称（用于 populate 函数路由）
    """
    if slide.layout_hint:
        return find_layout_flexible(prs, slide.layout_hint)

    has_chart = has_image = has_table = has_blockquote = has_text = False
    for e in slide.body_elements:
        if isinstance(e, ChartElement):
            has_chart = True
        elif isinstance(e, (ImageElement, MermaidElement)):
            has_image = True
        elif isinstance(e, TableElement):
            has_table = True
        elif isinstance(e, BlockquoteElement):
            has_blockquote = True
        elif isinstance(e, TextElement):
            has_text = True

    if has_chart and (has_text or has_table):
        std_name = 'mixed'
    elif has_chart:
        std_name = 'chart'
    elif has_image and not has_text and not has_table:
        std_name = 'image'
    elif has_blockquote or (has_text and has_image):
        std_name = 'two-column'
    elif has_table and not has_text:
        std_name = 'table'
    elif not slide.body_elements:
        std_name = 'title-only'
    else:
        std_name = 'standard'

    return find_layout_flexible(prs, std_name)


# ---------------------------------------------------------------------------
# Text formatting helpers
# ---------------------------------------------------------------------------

def set_text_frame(text_frame, paragraphs: list[Paragraph], clear: bool = True):
    """Fill a text frame with formatted paragraphs."""
    if clear:
        text_frame.clear()

    for pi, para in enumerate(paragraphs):
        if pi == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.level = para.level
        for ri, run_data in enumerate(para.runs):
            if ri == 0 and pi == 0 and clear:
                # Reuse the existing run in the first paragraph
                run = p.runs[0] if p.runs else p.add_run()
                run.text = run_data.text
            else:
                run = p.add_run()
                run.text = run_data.text
            if run_data.bold:
                run.font.bold = True
            if run_data.italic:
                run.font.italic = True


def collect_text_paragraphs(elements: list) -> list[Paragraph]:
    """Collect all text paragraphs from TextElement and BlockquoteElement lists."""
    paras = []
    for el in elements:
        if isinstance(el, TextElement):
            paras.extend(el.paragraphs)
        elif isinstance(el, BlockquoteElement):
            paras.extend(el.paragraphs)
    return paras


# ---------------------------------------------------------------------------
# Table helpers
# ---------------------------------------------------------------------------

def add_table_to_slide(slide, table_el: TableElement, left, top, width, height):
    """Add a formatted table shape to a slide."""
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    n_rows = len(table_el.rows) + 1  # +1 for header
    n_cols = len(table_el.headers)
    if n_cols == 0:
        return

    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = tbl_shape.table

    # Set column widths evenly
    col_width = width // n_cols
    for ci in range(n_cols):
        table.columns[ci].width = col_width

    # Fill header
    for ci, header in enumerate(table_el.headers):
        cell = table.cell(0, ci)
        cell.text = header
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.bold = True
                run.font.size = Pt(14)

    # Fill data rows
    for ri, row in enumerate(table_el.rows):
        for ci, cell_text in enumerate(row):
            if ci < n_cols:
                cell = table.cell(ri + 1, ci)
                cell.text = cell_text
                for para in cell.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.CENTER
                    for run in para.runs:
                        run.font.size = Pt(12)

    return tbl_shape


# ---------------------------------------------------------------------------
# Image helpers
# ---------------------------------------------------------------------------

def add_image_to_slide(slide, img_path: str, left, top, max_width, max_height):
    """Add an image to a slide, scaled to fit within bounds while preserving aspect ratio."""
    from PIL import Image

    if not os.path.exists(img_path):
        # Add a placeholder text box instead
        txBox = slide.shapes.add_textbox(left, top, max_width, Emu(400000))
        txBox.text_frame.text = f"[Image not found: {img_path}]"
        return txBox

    with Image.open(img_path) as img:
        img_w, img_h = img.size

    aspect = img_w / img_h
    target_w = max_width
    target_h = int(target_w / aspect)

    if target_h > max_height:
        target_h = max_height
        target_w = int(target_h * aspect)

    # Center within the available area
    x = left + (max_width - target_w) // 2
    y = top + (max_height - target_h) // 2

    return slide.shapes.add_picture(img_path, x, y, target_w, target_h)


# ---------------------------------------------------------------------------
# Mermaid rendering
# ---------------------------------------------------------------------------

def render_mermaid(code: str, index: int, img_dir: str = "slides-img") -> str:
    """Render a mermaid diagram to PNG, return the output path."""
    os.makedirs(img_dir, exist_ok=True)
    mmd_path = os.path.join(img_dir, f"mermaid-{index}.mmd")
    png_path = os.path.join(img_dir, f"mermaid-{index}.png")

    with open(mmd_path, 'w', encoding='utf-8') as f:
        f.write(code)

    try:
        subprocess.run(
            ["npx", "--yes", "@mermaid-js/mermaid-cli",
             "-i", mmd_path, "-o", png_path,
             "-b", "transparent", "-w", "1200", "--quiet"],
            check=True, capture_output=True, timeout=60
        )
    except (subprocess.CalledProcessError, FileNotFoundError) as e:
        print(f"  Warning: mermaid render failed for diagram {index}: {e}")
        return ""

    return png_path if os.path.exists(png_path) else ""


def render_all_mermaid(slides: list[SlideData], img_dir: str = "slides-img"):
    """Render all mermaid elements in slides, replacing them with ImageElements."""
    for slide in slides:
        new_elements = []
        for el in slide.body_elements:
            if isinstance(el, MermaidElement):
                png_path = render_mermaid(el.code, el.index, img_dir)
                if png_path:
                    new_elements.append(ImageElement(path=png_path, alt="mermaid diagram"))
                else:
                    new_elements.append(TextElement(
                        paragraphs=[Paragraph(runs=[Run(text="[Mermaid diagram render failed]")])]))
            else:
                new_elements.append(el)
        slide.body_elements = new_elements


# ---------------------------------------------------------------------------
# WPS compatibility fixup
# ---------------------------------------------------------------------------

def wps_fixup(pptx_path: str):
    """Post-process a saved pptx to ensure WPS compatibility.

    Checks:
    1. All parts in _rels files have corresponding Content_Types entries
    2. No orphaned relationship entries pointing to missing files
    """
    with zipfile.ZipFile(pptx_path, 'r') as zin:
        entries = {name: zin.read(name) for name in zin.namelist()}

    all_files = set(entries.keys())
    modified = False

    # Parse Content_Types
    ct_bytes = entries['[Content_Types].xml']
    ct = etree.fromstring(ct_bytes)
    ct_ns = 'http://schemas.openxmlformats.org/package/2006/content-types'

    # Map of part name -> content type for Override elements
    existing_overrides = set()
    for override in ct.findall(f'{{{ct_ns}}}Override'):
        existing_overrides.add(override.get('PartName'))

    # Content type mapping
    ct_map = {
        'slide': 'application/vnd.openxmlformats-officedocument.presentationml.slide+xml',
        'slideMaster': 'application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml',
        'slideLayout': 'application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml',
        'theme': 'application/vnd.openxmlformats-officedocument.theme+xml',
        'notesSlide': 'application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml',
        'chart': 'application/vnd.openxmlformats-officedocument.drawingml.chart+xml',
        'chartStyle': 'application/vnd.ms-office.chartstyle+xml',
        'chartColorStyle': 'application/vnd.ms-office.chartcolorstyle+xml',
    }

    # Check all slide-related XML files have Content_Types entries
    for fpath in all_files:
        if fpath.startswith('ppt/') and fpath.endswith('.xml'):
            part_name = '/' + fpath
            if part_name not in existing_overrides:
                for key, ctype in ct_map.items():
                    if key in fpath:
                        etree.SubElement(ct, f'{{{ct_ns}}}Override',
                                         PartName=part_name, ContentType=ctype)
                        existing_overrides.add(part_name)
                        modified = True
                        break

    # Check for orphaned rels (pointing to non-existent files)
    rel_ns = 'http://schemas.openxmlformats.org/package/2006/relationships'
    for fpath in list(all_files):
        if fpath.endswith('.rels'):
            rels_xml = etree.fromstring(entries[fpath])
            rels_dir = os.path.dirname(fpath.replace('_rels/', '').rstrip('/'))
            to_remove = []
            for rel in rels_xml:
                target = rel.get('Target', '')
                if target.startswith('http://') or target.startswith('https://'):
                    continue
                # Resolve relative path
                if target.startswith('/'):
                    resolved = target.lstrip('/')
                else:
                    resolved = os.path.normpath(os.path.join(rels_dir, target))
                if resolved not in all_files:
                    to_remove.append(rel)

            if to_remove:
                for rel in to_remove:
                    rels_xml.remove(rel)
                entries[fpath] = etree.tostring(rels_xml, xml_declaration=True,
                                                encoding='UTF-8', standalone=True)
                modified = True

    if modified:
        entries['[Content_Types].xml'] = etree.tostring(ct, xml_declaration=True,
                                                         encoding='UTF-8', standalone=True)
        with zipfile.ZipFile(pptx_path, 'w', zipfile.ZIP_DEFLATED) as zout:
            for name, data in entries.items():
                zout.writestr(name, data)


# ---------------------------------------------------------------------------
# Layout info extraction (for extract_template)
# ---------------------------------------------------------------------------

def get_layout_info(prs: Presentation) -> list[dict]:
    """Extract layout information from a presentation."""
    layouts = []
    for mi, master in enumerate(prs.slide_masters):
        for layout in master.slide_layouts:
            placeholders = []
            for ph in layout.placeholders:
                placeholders.append({
                    'idx': ph.placeholder_format.idx,
                    'name': ph.name,
                    'type': str(ph.placeholder_format.type),
                })
            layouts.append({
                'master': mi,
                'name': layout.name,
                'placeholders': placeholders,
            })
    return layouts


def get_used_layouts(prs: Presentation) -> list[dict]:
    """Get layouts actually used by slides in the presentation."""
    used = []
    seen = set()
    for slide in prs.slides:
        layout = slide.slide_layout
        for mi, master in enumerate(prs.slide_masters):
            if layout in master.slide_layouts:
                key = (mi, layout.name)
                if key not in seen:
                    seen.add(key)
                    used.append({
                        'master': mi,
                        'name': layout.name,
                        'placeholders': [
                            {'idx': p.placeholder_format.idx, 'name': p.name}
                            for p in layout.placeholders
                        ],
                    })
                break
    return used
